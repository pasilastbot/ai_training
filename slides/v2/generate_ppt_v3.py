#!/usr/bin/env python3
"""
Generate ModernPath AI Training PowerPoint v3 - Complete slide set with all imported slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches

def create_training_ppt_v3():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    base = "product/training-materials/v2"
    
    # All slides in logical order
    slides = [
        # ========== INTRO ==========
        ("00-master-title.png", "AI-NATIVE ENGINEERING MASTERY", 
         "Transform Your Development Practice\nTrainer: Pasi Vuorio | ModernPath"),
        ("01-intro-trainer.png", "Your Trainer: Pasi Vuorio",
         "Founder & AI Architect at ModernPath\n\n• 30+ years software engineering\n• Enterprise AI transformation specialist\n• Brownfield modernization expert"),
        ("02-training-overview.png", "Training Overview: Proto → Production → Scale",
         "SESSION 1: Fundamentals - Build rapid prototypes with AI\nSESSION 2: Engineering - Production-grade development\nSESSION 3: Scaling - Optimize & scale"),
        
        # ========== SESSION 1: FUNDAMENTALS (Proto) ==========
        ("01-session1-title.png", "SESSION 1: AI-Assisted Development Fundamentals",
         "90 min | Beginner\nGoal: Build rapid prototypes with AI"),
        ("02-s1-tools-landscape.png", "1. AI Coding Tools Landscape",
         "Cursor, Copilot, Claude Code, Codex"),
        ("03-s1-rules-config.png", "Rules & Configuration",
         ".cursor/rules/, CLAUDE.md, AGENTS.md - Your codebase constitution"),
        ("04-s1-context-sandwich.png", "2. Effective Prompting Patterns",
         "Context Sandwich: Evidence → Request → Constraints"),
        ("05-s1-model-selection.png", "3. Model Selection",
         "Claude vs GPT vs Gemini - Match the brain to the task"),
        ("06-s1-green-vs-brown.png", "4. Greenfield vs Brownfield",
         "80% of enterprise work is brownfield"),
        ("07-s1-greenfield-sdd.png", "5. Greenfield: Spec-Driven Development",
         "Write the contract before the code"),
        ("08-s1-brownfield-docs.png", "6. Brownfield: 3-Tiered Documentation",
         "Your docs are the map for the AI brain"),
        ("21-documentation-plan.png", "Documentation Plan: content-plan.md",
         "One index file for both greenfield and brownfield projects"),
        
        # ========== SESSION 2: ENGINEERING (Production) ==========
        ("09-session2-title.png", "SESSION 2: Spec-Driven Development & TDD",
         "90 min | Intermediate\nGoal: Production-grade development"),
        ("10-vibe-vs-sdd.png", "Vibe Coding vs Spec-Driven Development",
         "Vibe: Fast but fragile\nSDD: Slower start but maintainable"),
        ("10-s2-4step-process.png", "ModernPath 4-Step Process",
         "Document → Spec → Develop → Audit"),
        ("11-s2-specifications.png", "Writing Effective Specifications",
         "Templates, acceptance criteria, testable requirements"),
        ("12-s2-tdd-cycle.png", "TDD Cycle with AI",
         "Red → Green → Refactor"),
        ("13--architecture.png", "AI-Friendly Architecture",
         "Loose coupling = safe AI edits"),
        ("13-checkpoints.png", "Quality Checkpoints",
         "Gates to ensure quality at each stage"),
        ("22-s2-code-review.png", "Code Review with AI",
         "Critique loops, PR review prompts, quality feedback"),
        ("23-s2-fix-workflow.png", "The Fix Workflow",
         "Reproduce → Hypothesize → Validate → Fix → Document"),
        
        # ========== SESSION 3: SCALING ==========
        ("14-session3-title.png", "SESSION 3: Advanced Agentic Workflows",
         "90 min | Advanced\nGoal: Optimize & scale"),
        ("15-s3-agentic-architecture.png", "Agentic Architecture",
         "Assistant → Tool User → Agent → Multi-Agent"),
        ("16-s3-sub-agents.png", "Specialized Sub-Agents",
         "Spec writer, test generator, code reviewer"),
        ("17-s3-mcp-vs-cli.png", "MCP vs CLI Tools",
         "CLI: 30%+ better tokens, composable, debuggable, but can hallucinate syntax.\nMCP: Type-safe, complex integrations, approvals, but context bloat.\n\nVerdict: CLI outperforms for existing tools. MCP for specialized/secure workflows."),
        ("18-s3-mcp-integrations.png", "MCP Integrations",
         "GitHub, Slack, databases, custom servers"),
        ("19-s3-hooks-verification.png", "Hooks & Verification",
         "PostToolUse, Stop hooks, permission management"),
        ("20-s3-scaling.png", "Scaling Across Teams",
         "Shared agents, centralized docs, metrics"),
        ("24-s3-agent-skills.png", "Agent Skills",
         "Web search, file operations, code analysis, external APIs"),
        ("25-s3-inventing-cli.png", "Inventing CLI Tools",
         "Need → Script → Test → Register"),
        ("26-s3-skills-agents-hooks.png", "Skills vs Sub-Agents vs Hooks",
         "SKILLS = WHAT (capabilities)\nSUB-AGENTS = WHO (specialists)\nHOOKS = WHEN (control)"),
    ]
    
    # Generate presentation
    count = 0
    missing = []
    for filename, title, notes in slides:
        path = f"{base}/{filename}"
        if not os.path.exists(path):
            missing.append(filename)
            continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"{title}\n\n{notes}"
        
        print(f"✓ {title}")
        count += 1
    
    if missing:
        print(f"\n⚠️  Missing files: {missing}")
    
    output = f"{base}/ModernPath_AI_Training_v3.pptx"
    prs.save(output)
    print(f"\n✅ Saved: {output}")
    print(f"📊 Total slides: {count}")

if __name__ == "__main__":
    create_training_ppt_v3()
