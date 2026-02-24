#!/usr/bin/env python3
"""
Generate ModernPath AI Training PowerPoint v2 - All new professional slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches

def create_training_ppt_v2():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    base = "product/training-materials/v2"
    
    # All slides in order
    slides = [
        ("00-master-title.png", "AI-NATIVE ENGINEERING MASTERY", 
         "Transform Your Development Practice\nTrainer: Pasi Vuorio | ModernPath"),
        ("01-intro-trainer.png", "Your Trainer: Pasi Vuorio",
         "Founder & AI Architect at ModernPath\n\n• 30+ years software engineering\n• Enterprise AI transformation specialist\n• Brownfield modernization expert\n\nModernPath: AI-Native Software Development"),
        ("02-training-overview.png", "Training Overview: Proto → Production → Scale",
         "SESSION 1: Fundamentals - Build rapid prototypes with AI\nSESSION 2: Engineering - Production-grade development (specs, TDD, architecture)\nSESSION 3: Scaling - Optimize & scale (multi-agent, teams, metrics)"),
        
        # SESSION 1
        ("01-session1-title.png", "SESSION 1: AI-Assisted Development Fundamentals",
         "90 min | Beginner"),
        ("02-s1-tools-landscape.png", "1. AI Coding Tools Landscape",
         "Cursor, Copilot, Claude Code, Codex"),
        ("03-s1-rules-config.png", "Rules & Configuration",
         ".cursor/rules/, CLAUDE.md, AGENTS.md - Your codebase constitution"),
        ("04-s1-context-sandwich.png", "2. Effective Prompting Patterns",
         "Context Sandwich: Evidence → Request → Constraints"),
        ("05-s1-model-selection.png", "3. Model Selection",
         "Claude vs GPT vs Gemini - Match the brain to the task"),
        ("06-s1-green-vs-brown.png", "4. Greenfield vs Brownfield",
         "80% of enterprise work is brownfield"),
        ("07-s1-greenfield-sdd.png", "5. Greenfield: Spec-Driven Development",
         "Write the contract before the code"),
        ("08-s1-brownfield-docs.png", "6. Brownfield: 3-Tiered Documentation",
         "Your docs are the map for the AI brain"),
        
        # SESSION 2
        ("09-session2-title.png", "SESSION 2: Spec-Driven Development & TDD",
         "90 min | Intermediate"),
        ("10-s2-4step-process.png", "ModernPath 4-Step Process",
         "Document → Spec → Develop → Audit"),
        ("11-s2-specifications.png", "Writing Effective Specifications",
         "Templates, acceptance criteria, testable requirements"),
        ("12-s2-tdd-cycle.png", "TDD Cycle with AI",
         "Red → Green → Refactor"),
        ("13-s2-architecture.png", "AI-Friendly Architecture",
         "Loose coupling = safe AI edits"),
        ("22-s2-code-review.png", "Code Review with AI",
         "Critique loops, PR review prompts, quality feedback.\n\nAI catches what humans miss at 3am."),
        ("23-s2-fix-workflow.png", "The Fix Workflow",
         "Systematic debugging: Reproduce → Hypothesize → Validate → Fix → Document.\n\nEvery fix becomes team knowledge in learnings.md."),
        
        # SESSION 3
        ("14-session3-title.png", "SESSION 3: Advanced Agentic Workflows",
         "90 min | Advanced"),
        ("15-s3-agentic-architecture.png", "Agentic Architecture",
         "Assistant → Tool User → Agent → Multi-Agent"),
        ("16-s3-sub-agents.png", "Specialized Sub-Agents",
         "Spec writer, test generator, code reviewer"),
        ("17-s3-mcp-vs-cli.png", "MCP vs CLI Tools",
         "MCP for efficiency, CLI for debugging"),
        ("18-s3-mcp-integrations.png", "MCP Integrations",
         "GitHub, Slack, databases, custom servers"),
        ("19-s3-hooks-verification.png", "Hooks & Verification",
         "PostToolUse, Stop hooks, permission management"),
        ("20-s3-scaling.png", "Scaling Across Teams",
         "Shared agents, centralized docs, metrics"),
        ("24-s3-agent-skills.png", "Agent Skills",
         "Web search, file operations, code analysis, external APIs.\n\nSkills = Superpowers for your agents."),
        ("25-s3-inventing-cli.png", "Inventing CLI Tools",
         "Need → Script → Test → Register.\n\nYour custom tools become AI superpowers."),
        ("26-s3-skills-agents-hooks.png", "Skills vs Sub-Agents vs Hooks",
         "SKILLS = WHAT (capabilities) - Add abilities like search, APIs\nSUB-AGENTS = WHO (specialists) - Delegate to experts\nHOOKS = WHEN (control) - Verify before/after actions"),
        
        # BONUS: Documentation Plan (applies to both greenfield & brownfield)
        ("21-documentation-plan.png", "Documentation Plan: content-plan.md",
         "One index file for both greenfield and brownfield projects.\n\nGreenfield: Plan before you build, define modules upfront\nBrownfield: Map what exists, index for AI navigation"),
    ]
    
    # Generate presentation
    count = 0
    for filename, title, notes in slides:
        path = f"{base}/{filename}"
        if not os.path.exists(path):
            print(f"⚠️  Missing: {path}")
            continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"{title}\n\n{notes}"
        
        print(f"✓ {title}")
        count += 1
    
    output = f"{base}/ModernPath_AI_Training_v2.pptx"
    prs.save(output)
    print(f"\n✅ Saved: {output}")
    print(f"📊 Total slides: {count}")

if __name__ == "__main__":
    create_training_ppt_v2()
